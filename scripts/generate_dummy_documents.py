"""
Generate dummy enterprise documents for the Enterprise RAG Assistant.

Creates:
  - HR_Policy.pdf            (PDF  - HR policies)
  - Leave_Policy.pptx        (PPTX - Leave policies)
  - Finance_Policy.docx      (DOCX - Finance policies)
  - Termination_Policy_Scanned.pdf (Scanned PDF - Termination)
  - Employee_Salary_Records.xlsx   (XLSX - Salary records)
"""

from pathlib import Path
import sys

# Add parent to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

# Ensure output directory
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "data" / "documents"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def generate_hr_policy_pdf():
    """Generate HR Policy PDF using fpdf or reportlab, or a simple text-based method."""
    print("  Creating HR_Policy.pdf...")
    try:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()

        # Title
        pdf.set_font("Arial", "B", 18)
        pdf.cell(0, 15, "HR Policy Document", ln=True, align="C")
        pdf.ln(10)

        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "1. Company Vision and Values", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "Our company is committed to fostering an inclusive, innovative, and "
            "high-performance work environment. We value integrity, teamwork, and "
            "continuous improvement in everything we do.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "2. Working Hours", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "Standard working hours are from 9:30 AM to 6:30 PM, Monday through Friday. "
            "Employees are expected to be at their workstations during these hours. "
            "Flexible timing options are available with manager approval.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "3. Probation Period", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "All new employees undergo a probation period of 6 months from the date of joining. "
            "Performance reviews are conducted at the end of the 3rd and 6th month. "
            "Confirmation of employment is subject to satisfactory performance.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "4. Attendance Policy", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "Employees must mark their attendance daily through the biometric system. "
            "Late arrival beyond 15 minutes is considered half-day leave. "
            "Regular attendance is rewarded through the attendance bonus program.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "5. Dress Code", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "Employees are expected to maintain a formal dress code during working hours. "
            "Business casuals are permitted on Fridays. "
            "On Fridays, casual dress is allowed. "
            "Special cultural events may have specific dress requirements.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "6. Code of Conduct", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "All employees must adhere to the company's code of conduct, which includes "
            "maintaining confidentiality, avoiding conflicts of interest, and treating "
            "colleagues with respect and dignity at all times.")

        pdf.ln(5)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "7. Grievance Process", ln=True)
        pdf.set_font("Arial", "", 11)
        pdf.multi_cell(0, 6,
            "Employees can raise grievances through the HR department or via the company "
            "grievance portal. Complaints are addressed within 7 business days. "
            "All grievances are handled confidentially and without fear of retaliation.")

        pdf.output(str(OUTPUT_DIR / "HR_Policy.pdf"))
        print("    Done.")
    except ImportError:
        print("    fpdf not available, creating using reportlab...")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas

            c = canvas.Canvas(str(OUTPUT_DIR / "HR_Policy.pdf"), pagesize=A4)
            c.setFont("Helvetica-Bold", 18)
            c.drawString(50, 800, "HR Policy Document")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 770, "1. Company Vision and Values")
            c.setFont("Helvetica", 11)
            c.drawString(50, 750, "Our company is committed to fostering an inclusive, innovative, and")
            c.drawString(50, 735, "high-performance work environment. We value integrity, teamwork, and")
            c.drawString(50, 720, "continuous improvement in everything we do.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 695, "2. Working Hours")
            c.setFont("Helvetica", 11)
            c.drawString(50, 675, "Standard working hours are from 9:30 AM to 6:30 PM, Monday through")
            c.drawString(50, 660, "Friday. Employees are expected to be at their workstations during")
            c.drawString(50, 645, "these hours.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 620, "3. Probation Period")
            c.setFont("Helvetica", 11)
            c.drawString(50, 600, "All new employees undergo a probation period of 6 months from the")
            c.drawString(50, 585, "date of joining. Performance reviews are conducted at the end of")
            c.drawString(50, 570, "the 3rd and 6th month.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 545, "4. Attendance Policy")
            c.setFont("Helvetica", 11)
            c.drawString(50, 525, "Employees must mark their attendance daily through the biometric")
            c.drawString(50, 510, "system. Late arrival beyond 15 minutes is considered half-day leave.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 485, "5. Dress Code")
            c.setFont("Helvetica", 11)
            c.drawString(50, 465, "Employees are expected to maintain a formal dress code during")
            c.drawString(50, 450, "working hours. Business casuals are permitted on Fridays.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 425, "6. Code of Conduct")
            c.setFont("Helvetica", 11)
            c.drawString(50, 405, "All employees must adhere to the company's code of conduct.")
            c.drawString(50, 390, "Maintain confidentiality and treat colleagues with respect.")
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 365, "7. Grievance Process")
            c.setFont("Helvetica", 11)
            c.drawString(50, 345, "Employees can raise grievances through the HR department or via")
            c.drawString(50, 330, "the company grievance portal. Complaints are addressed within")
            c.drawString(50, 315, "7 business days.")
            c.save()
            print("    Done.")
        except ImportError:
            print("    PDF libraries not available. Creating text file instead.")
            # Fallback: create a text file
            with open(OUTPUT_DIR / "HR_Policy.pdf", "w") as f:
                f.write("HR Policy Document\n\n")
                f.write("1. Company Vision and Values\n")
                f.write("...\n")


def generate_leave_policy_pptx():
    """Generate Leave Policy PowerPoint."""
    print("  Creating Leave_Policy.pptx...")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Leave Policy"
        subtitle.text = "Enterprise Leave Management - 2026"

        # Slide 2: Leave Overview
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Leave Overview"
        content = slide.placeholders[1]
        content.text = ("Our company offers a comprehensive leave policy designed to support "
                        "employee wellbeing and work-life balance.\n\n"
                        "Leave Types:\n"
                        "• Casual Leave: 12 days\n"
                        "• Sick Leave: 10 days\n"
                        "• Earned Leave: 18 days")

        # Slide 3: Casual Leave
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Casual Leave Policy"
        content = slide.placeholders[1]
        content.text = ("Casual Leave: 12 days per year\n\n"
                        "• For urgent personal matters\n"
                        "• Minimum 1 day, maximum 3 consecutive days\n"
                        "• Requires manager approval\n"
                        "• Cannot be clubbed with other leave types\n"
                        "• Unused casual leaves lapse at year end")

        # Slide 4: Sick Leave
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Sick Leave Policy"
        content = slide.placeholders[1]
        content.text = ("Sick Leave: 10 days per year\n\n"
                        "• For medical emergencies and illness\n"
                        "• Medical certificate required for 3+ consecutive days\n"
                        "• Unused sick leaves lapse at year end\n"
                        "• Can be used for immediate family medical needs")

        # Slide 5: Earned Leave
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Earned Leave Policy"
        content = slide.placeholders[1]
        content.text = ("Earned Leave: 18 days per year\n\n"
                        "• Accrues monthly (1.5 days per month)\n"
                        "• Maximum accumulation: 45 days\n"
                        "• Requires 2 weeks advance notice\n"
                        "• Encashment allowed up to 30 days at retirement/separation\n"
                        "• Can be clubbed with other leave types")

        # Slide 6: Leave Approval Workflow
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Leave Approval Workflow"
        content = slide.placeholders[1]
        content.text = ("Leave Request → Reporting Manager → HR Approval\n\n"
                        "• 1-3 days: Manager approval only\n"
                        "• 4-10 days: Manager + HR approval\n"
                        "• 10+ days: Manager + HR + Director approval\n"
                        "• Emergency leave: Inform manager within 24 hours\n"
                        "• All leaves must be applied through the HRMS portal")

        # Slide 7: Holiday Calendar
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Holiday Calendar"
        content = slide.placeholders[1]
        content.text = ("Company Holidays (2026):\n\n"
                        "• Republic Day: January 26\n"
                        "• Holi: March 14\n"
                        "• Independence Day: August 15\n"
                        "• Diwali: November 1 (plus optional day)\n"
                        "• Christmas: December 25\n\n"
                        "Total: 10 fixed holidays + 2 floating holidays")

        # Slide 8: FAQs
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Leave FAQs"
        content = slide.placeholders[1]
        content.text = ("Q: Can I carry forward unused leaves?\n"
                        "A: Earned Leave can be carried forward up to 45 days.\n\n"
                        "Q: What is the notice period for leave?\n"
                        "A: 2 weeks for Earned Leave, 1 day for Casual/Sick Leave.\n\n"
                        "Q: Can I encash my leaves?\n"
                        "A: Earned Leave encashment is allowed up to 30 days.\n\n"
                        "Q: What happens to leaves during notice period?\n"
                        "A: Unused Earned Leave is encashed in full and final settlement.")

        prs.save(str(OUTPUT_DIR / "Leave_Policy.pptx"))
        print("    Done.")
    except ImportError:
        print("    python-pptx not available. Skipping PPTX creation.")


def generate_finance_policy_docx():
    """Generate Finance Policy Word document."""
    print("  Creating Finance_Policy.docx...")
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        title = doc.add_heading("Finance Policy Document", level=0)
        doc.add_paragraph("Enterprise Finance & Reimbursement Guidelines - 2026")

        # Section 1
        doc.add_heading("1. Travel Reimbursement", level=1)
        doc.add_paragraph(
            "Employees are eligible for travel reimbursement for business-related travel. "
            "The following guidelines apply:"
        )
        doc.add_paragraph("Economy class airfare for domestic travel", style='List Bullet')
        doc.add_paragraph("Business class airfare for international travel over 6 hours", style='List Bullet')
        doc.add_paragraph("Train travel: AC 2-tier or above", style='List Bullet')
        doc.add_paragraph("Cab fare: Actuals with proper receipts", style='List Bullet')
        doc.add_paragraph("Personal vehicle: ₹12 per km", style='List Bullet')

        # Section 2
        doc.add_heading("2. Meal Allowance", level=1)
        doc.add_paragraph(
            "Employees traveling on business are entitled to a daily meal allowance. "
            "The current meal allowance is ₹800 per day. "
            "This covers breakfast, lunch, and dinner expenses during business travel. "
            "No additional documentation is required for meals within the allowance limit."
        )

        # Section 3
        doc.add_heading("3. Hotel Reimbursement", level=1)
        doc.add_paragraph(
            "Hotel accommodation during business travel is reimbursed as per the following limits:"
        )
        doc.add_paragraph("Tier 1 cities (Mumbai, Delhi, Bangalore): Up to ₹5,000 per night", style='List Bullet')
        doc.add_paragraph("Tier 2 cities: Up to ₹3,000 per night", style='List Bullet')
        doc.add_paragraph("Tier 3 cities: Up to ₹1,500 per night", style='List Bullet')

        # Section 4
        doc.add_heading("4. Internet Reimbursement", level=1)
        doc.add_paragraph(
            "Employees working from home are eligible for internet reimbursement "
            "of up to ₹1,500 per month. "
            "This requires submission of the internet bill for the respective month. "
            "The reimbursement is processed along with the monthly salary."
        )

        # Section 5
        doc.add_heading("5. Expense Claim Process", level=1)
        doc.add_paragraph(
            "All expense claims must be submitted through the company's expense management system "
            "within 30 days of the expense being incurred. Claims require:"
        )
        doc.add_paragraph("Itemized bills/receipts for all expenses", style='List Bullet')
        doc.add_paragraph("Manager approval for pre-travel expenses", style='List Bullet')
        doc.add_paragraph("Business purpose description", style='List Bullet')
        doc.add_paragraph("Original receipts for expenses above ₹1,000", style='List Bullet')

        # Section 6
        doc.add_heading("6. Payment Timeline", level=1)
        doc.add_paragraph(
            "Approved expense claims are processed within the following timelines:"
        )
        doc.add_paragraph("Standard claims: Processed within 15 business days", style='List Bullet')
        doc.add_paragraph("Urgent claims (marked as urgent): Processed within 5 business days", style='List Bullet')
        doc.add_paragraph("Advance requests: Processed within 3 business days before travel", style='List Bullet')

        doc.save(str(OUTPUT_DIR / "Finance_Policy.docx"))
        print("    Done.")
    except ImportError:
        print("    python-docx not available. Skipping DOCX creation.")


def generate_employee_salary_xlsx():
    """Generate Employee Salary Records Excel file."""
    print("  Creating Employee_Salary_Records.xlsx...")
    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Employee Salary Records"

        # Define styles
        header_font = Font(bold=True, color="FFFFFF", size=12)
        header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin'),
        )

        # Headers
        headers = ["Employee ID", "Name", "Department", "Designation",
                    "Monthly Salary", "Annual CTC", "Date of Joining", "Bank Account"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = thin_border

        # Employee data
        employees = [
            ["EMP1001", "Rahul Sharma", "Software", "Software Engineer",
             "₹75,000", "10.2 LPA", "2024-01-15", "XXXX-XXXX-1234"],
            ["EMP1002", "Sneha Verma", "Software", "Senior Software Engineer",
             "₹98,000", "13.4 LPA", "2023-06-01", "XXXX-XXXX-5678"],
            ["EMP1003", "Amit Singh", "HR", "HR Executive",
             "₹62,000", "8.5 LPA", "2024-03-10", "XXXX-XXXX-9012"],
            ["EMP1004", "Neha Kapoor", "Finance", "Finance Analyst",
             "₹82,000", "11.0 LPA", "2023-11-20", "XXXX-XXXX-3456"],
            ["EMP1005", "Vikram Joshi", "Sales", "Sales Executive",
             "₹68,000", "9.2 LPA", "2024-07-01", "XXXX-XXXX-7890"],
        ]

        data_font = Font(size=11)
        data_alignment = Alignment(horizontal="center", vertical="center")

        for row_idx, emp in enumerate(employees, 2):
            for col_idx, value in enumerate(emp, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.font = data_font
                cell.alignment = data_alignment
                cell.border = thin_border

        # Adjust column widths
        column_widths = [15, 18, 15, 25, 18, 15, 18, 20]
        for i, width in enumerate(column_widths, 1):
            ws.column_dimensions[chr(64 + i)].width = width

        wb.save(str(OUTPUT_DIR / "Employee_Salary_Records.xlsx"))
        print("    Done.")
    except ImportError:
        print("    openpyxl not available. Skipping XLSX creation.")


def generate_termination_policy_scanned_pdf():
    """Generate a termination policy PDF that simulates a scanned document."""
    print("  Creating Termination_Policy_Scanned.pdf...")
    try:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()

        # Simulate a scanned document look
        pdf.set_font("Courier", "B", 16)
        pdf.cell(0, 15, "TERMINATION POLICY", ln=True, align="C")
        pdf.ln(5)

        pdf.set_font("Courier", "", 12)
        pdf.multi_cell(0, 7,
            "This document outlines the resignation and termination"
            " procedures for all employees."
        )
        pdf.ln(5)

        pdf.set_font("Courier", "B", 12)
        pdf.cell(0, 10, "1. Resignation Procedure", ln=True)
        pdf.set_font("Courier", "", 11)
        pdf.multi_cell(0, 7,
            "Employees wishing to resign must submit a formal resignation"
            " letter to their reporting manager. The notice period begins"
            " upon acceptance of the resignation by the manager."
        )

        pdf.ln(3)
        pdf.set_font("Courier", "B", 12)
        pdf.cell(0, 10, "2. Notice Period", ln=True)
        pdf.set_font("Courier", "", 11)
        pdf.multi_cell(0, 7,
            "Notice Period Details:\n"
            "  - Permanent Employees: 60 days\n"
            "  - Probation Employees: 15 days\n"
            "  - During Probation: 15 days\n\n"
            "The notice period may be waived or reduced at the company's"
            " discretion. Garden leave may be applied during the notice"
            " period at management's discretion."
        )

        pdf.ln(3)
        pdf.set_font("Courier", "B", 12)
        pdf.cell(0, 10, "3. Exit Clearance", ln=True)
        pdf.set_font("Courier", "", 11)
        pdf.multi_cell(0, 7,
            "All employees must complete the exit clearance process before"
            " final settlement. This includes:\n"
            "  - Return of company property (laptop, ID card, access cards)\n"
            "  - Clearance from IT department\n"
            "  - Clearance from Finance department\n"
            "  - Clearance from Administration\n"
            "  - Handover of pending work"
        )

        pdf.ln(3)
        pdf.set_font("Courier", "B", 12)
        pdf.cell(0, 10, "4. Full and Final Settlement", ln=True)
        pdf.set_font("Courier", "", 11)
        pdf.multi_cell(0, 7,
            "The full and final settlement will be processed within 30 days"
            " from the last working day. Settlement includes:\n"
            "  - Salary for days worked\n"
            "  - Encashment of earned leaves\n"
            "  - Any other dues as per company policy\n\n"
            "Any deductions will be applied as per company policy and"
            " applicable laws."
        )

        pdf.ln(5)
        pdf.set_font("Courier", "", 10)
        pdf.multi_cell(0, 7,
            "________________________________\n"
            "HR Manager Signature & Date\n\n"
            "________________________________\n"
            "Company Stamp"
        )

        pdf.output(str(OUTPUT_DIR / "Termination_Policy_Scanned.pdf"))
        print("    Done.")
    except ImportError:
        print("    fpdf not available. Using reportlab...")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            c = canvas.Canvas(str(OUTPUT_DIR / "Termination_Policy_Scanned.pdf"), pagesize=A4)
            c.setFont("Courier-Bold", 16)
            c.drawString(50, 800, "TERMINATION POLICY")
            c.setFont("Courier", 12)
            c.drawString(50, 770, "This document outlines the resignation and")
            c.drawString(50, 755, "termination procedures for all employees.")
            c.setFont("Courier-Bold", 12)
            c.drawString(50, 730, "1. Resignation Procedure")
            c.setFont("Courier", 11)
            c.drawString(50, 710, "Employees wishing to resign must submit a formal")
            c.drawString(50, 695, "resignation letter to their reporting manager.")
            c.setFont("Courier-Bold", 12)
            c.drawString(50, 670, "2. Notice Period")
            c.setFont("Courier", 11)
            c.drawString(50, 650, "- Permanent Employees: 60 days")
            c.drawString(50, 635, "- Probation Employees: 15 days")
            c.setFont("Courier-Bold", 12)
            c.drawString(50, 610, "3. Exit Clearance")
            c.setFont("Courier", 11)
            c.drawString(50, 590, "Return company property and complete clearance")
            c.drawString(50, 575, "from IT, Finance, and Administration.")
            c.setFont("Courier-Bold", 12)
            c.drawString(50, 550, "4. Full and Final Settlement")
            c.setFont("Courier", 11)
            c.drawString(50, 530, "Processed within 30 days from last working day.")
            c.drawString(50, 510, "Includes salary, leave encashment, and other dues.")
            c.setFont("Courier", 10)
            c.drawString(50, 480, "________________________________")
            c.drawString(50, 465, "HR Manager Signature & Date")
            c.drawString(50, 440, "________________________________")
            c.drawString(50, 425, "Company Stamp")
            c.save()
            print("    Done.")
        except ImportError:
            print("    PDF libraries not available. Creating text placeholder.")
            with open(OUTPUT_DIR / "Termination_Policy_Scanned.pdf", "w") as f:
                f.write("TERMINATION POLICY (Scanned Document Simulation)\n")


def main():
    print("=" * 60)
    print("  Generating Enterprise Dummy Documents")
    print("=" * 60)
    print(f"  Output directory: {OUTPUT_DIR}\n")

    generate_hr_policy_pdf()
    generate_leave_policy_pptx()
    generate_finance_policy_docx()
    generate_termination_policy_scanned_pdf()
    generate_employee_salary_xlsx()

    print("\n" + "=" * 60)
    print("  Document generation complete!")
    print("=" * 60)
    print("\n  Generated files:")
    for f in sorted(OUTPUT_DIR.iterdir()):
        size = f.stat().st_size
        print(f"    • {f.name} ({size:,} bytes)")
    print()


if __name__ == "__main__":
    main()