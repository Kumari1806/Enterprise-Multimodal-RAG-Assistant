"""Document ingestion module — multi-format parsing, OCR, chunking, embedding."""

import logging
from pathlib import Path
from typing import Optional

import chromadb
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
import pandas as pd
import numpy as np

from backend.config import (
    DOCUMENTS_DIR, CHROMA_DIR, CHROMA_COLLECTION_NAME,
    CHUNK_SIZE, CHUNK_OVERLAP,
    GEMINI_EMBEDDING_MODEL, GOOGLE_API_KEY,
    ROLE_DOCUMENTS, DOCUMENT_METADATA,
)

logger = logging.getLogger(__name__)


class DocumentIngestor:
    """Handles ingestion of enterprise documents into ChromaDB."""

    def __init__(self):
        self.chroma_client = chromadb.PersistentClient(path=str(CHROMA_DIR))
        self.embeddings = None
        if GOOGLE_API_KEY:
            try:
                self.embeddings = GoogleGenerativeAIEmbeddings(
                    model=GEMINI_EMBEDDING_MODEL,
                    google_api_key=GOOGLE_API_KEY,
                )
            except Exception as e:
                logger.warning(f"Could not initialize Gemini embeddings: {e}")
                self.embeddings = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

    def _get_or_create_collection(self):
        """Get or create the ChromaDB collection."""
        try:
            return self.chroma_client.get_collection(CHROMA_COLLECTION_NAME)
        except Exception:
            return self.chroma_client.create_collection(CHROMA_COLLECTION_NAME)

    def extract_text_from_pdf(self, filepath: Path) -> list[dict]:
        """Extract text from a PDF file with page numbers."""
        pages = []
        try:
            reader = PdfReader(str(filepath))
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append({
                        "content": text.strip(),
                        "metadata": {
                            "document": filepath.name,
                            "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                            "page": i,
                            "format": "PDF",
                        }
                    })
            logger.info(f"Extracted {len(pages)} pages from PDF: {filepath.name}")
        except Exception as e:
            logger.error(f"Error extracting PDF {filepath.name}: {e}")
        return pages

    def extract_text_from_pptx(self, filepath: Path) -> list[dict]:
        """Extract text from a PowerPoint file with slide numbers."""
        slides = []
        try:
            prs = Presentation(str(filepath))
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append({
                        "content": "\n".join(slide_text),
                        "metadata": {
                            "document": filepath.name,
                            "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                            "slide": i,
                            "format": "PowerPoint",
                        }
                    })
            logger.info(f"Extracted {len(slides)} slides from PPTX: {filepath.name}")
        except Exception as e:
            logger.error(f"Error extracting PPTX {filepath.name}: {e}")
        return slides

    def extract_text_from_docx(self, filepath: Path) -> list[dict]:
        """Extract text from a Word document."""
        sections = []
        try:
            doc = DocxDocument(str(filepath))
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text.strip())
            if full_text:
                sections.append({
                    "content": "\n".join(full_text),
                    "metadata": {
                        "document": filepath.name,
                        "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                        "format": "Word",
                        "section": "main",
                    }
                })
            # Also extract tables
            for i, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    sections.append({
                        "content": "\n".join(table_text),
                        "metadata": {
                            "document": filepath.name,
                            "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                            "format": "Word",
                            "section": f"table_{i+1}",
                        }
                    })
            logger.info(f"Extracted {len(sections)} sections from DOCX: {filepath.name}")
        except Exception as e:
            logger.error(f"Error extracting DOCX {filepath.name}: {e}")
        return sections

    def extract_text_from_xlsx(self, filepath: Path) -> list[dict]:
        """Extract text from an Excel file."""
        sheets_data = []
        try:
            xls = pd.ExcelFile(str(filepath))
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(str(filepath), sheet_name=sheet_name)
                # Convert to text representation
                text_lines = []
                # Add headers
                headers = [str(col) for col in df.columns]
                text_lines.append(" | ".join(headers))
                # Add rows
                for _, row in df.iterrows():
                    row_text = [str(val) if pd.notna(val) else "" for val in row]
                    text_lines.append(" | ".join(row_text))
                if text_lines:
                    sheets_data.append({
                        "content": "\n".join(text_lines),
                        "metadata": {
                            "document": filepath.name,
                            "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                            "format": "Excel",
                            "sheet": sheet_name,
                        }
                    })
            logger.info(f"Extracted {len(sheets_data)} sheets from XLSX: {filepath.name}")
        except Exception as e:
            logger.error(f"Error extracting XLSX {filepath.name}: {e}")
        return sheets_data

    def extract_text_from_scanned_pdf(self, filepath: Path) -> list[dict]:
        """Extract text from a scanned PDF using OCR (EasyOCR)."""
        pages = []
        try:
            import fitz  # PyMuPDF
            import easyocr

            reader = easyocr.Reader(["en"], gpu=False)
            doc = fitz.open(str(filepath))

            for i in range(len(doc)):
                page = doc[i]
                # Render page to image
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")

                # Save temp image for OCR
                temp_img_path = Path("/tmp") / f"ocr_page_{i}.png"
                with open(temp_img_path, "wb") as f:
                    f.write(img_data)

                # OCR
                try:
                    results = reader.readtext(str(temp_img_path))
                    text = " ".join([res[1] for res in results])
                    if text.strip():
                        pages.append({
                            "content": text.strip(),
                            "metadata": {
                                "document": filepath.name,
                                "document_name": DOCUMENT_METADATA.get(filepath.name, {}).get("name", filepath.name),
                                "page": i + 1,
                                "format": "Scanned PDF",
                                "ocr": True,
                            }
                        })
                finally:
                    if temp_img_path.exists():
                        temp_img_path.unlink()

                doc.close()
            logger.info(f"Extracted {len(pages)} pages via OCR from: {filepath.name}")
        except ImportError as e:
            logger.error(f"OCR libraries not available: {e}. Cannot process scanned PDF.")
        except Exception as e:
            logger.error(f"Error during OCR for {filepath.name}: {e}")
        return pages

    def extract_text(self, filepath: Path) -> list[dict]:
        """Extract text from a document based on its extension."""
        ext = filepath.suffix.lower()
        if ext == ".pdf":
            # Check if it's a scanned PDF (try regular extraction first)
            pages = self.extract_text_from_pdf(filepath)
            if not pages or all(len(p["content"]) < 50 for p in pages):
                # Probably scanned, try OCR
                logger.info(f"PDF {filepath.name} appears scanned, attempting OCR...")
                ocr_pages = self.extract_text_from_scanned_pdf(filepath)
                if ocr_pages:
                    return ocr_pages
            return pages
        elif ext == ".pptx":
            return self.extract_text_from_pptx(filepath)
        elif ext == ".docx":
            return self.extract_text_from_docx(filepath)
        elif ext == ".xlsx":
            return self.extract_text_from_xlsx(filepath)
        else:
            logger.warning(f"Unsupported file format: {ext}")
            return []

    def chunk_and_embed(self, extracted_data: list[dict], filename: str) -> int:
        """Chunk extracted text and index into ChromaDB."""
        collection = self._get_or_create_collection()
        total_chunks = 0

        for item in extracted_data:
            content = item["content"]
            metadata = item["metadata"]

            # Chunk the content
            chunks = self.text_splitter.split_text(content)

            for chunk_idx, chunk in enumerate(chunks):
                if not chunk.strip():
                    continue

                # Add role document mapping
                chunk_metadata = metadata.copy()
                chunk_metadata["chunk_index"] = chunk_idx
                chunk_metadata["filename"] = filename

                # Add role permissions
                roles_with_access = []
                for role, docs in ROLE_DOCUMENTS.items():
                    if filename in docs:
                        roles_with_access.append(role)
                chunk_metadata["roles_with_access"] = ",".join(roles_with_access)

                # Generate a unique ID
                chunk_id = f"{filename}_chunk_{total_chunks}"

                try:
                    if self.embeddings:
                        embedding = self.embeddings.embed_documents([chunk])[0]
                        collection.add(
                            embeddings=[embedding],
                            documents=[chunk],
                            metadatas=[chunk_metadata],
                            ids=[chunk_id],
                        )
                    else:
                        # Without Gemini, let ChromaDB use its built-in embedding
                        collection.add(
                            documents=[chunk],
                            metadatas=[chunk_metadata],
                            ids=[chunk_id],
                        )
                except Exception as e:
                    logger.warning(f"Error adding chunk {chunk_id}: {e}")
                    continue

                total_chunks += 1

        return total_chunks

    def ingest_all_documents(self) -> dict[str, int]:
        """Ingest all documents from the documents directory."""
        results = {}
        for filepath in sorted(DOCUMENTS_DIR.iterdir()):
            if filepath.is_file() and filepath.suffix.lower() in [".pdf", ".pptx", ".docx", ".xlsx"]:
                logger.info(f"Ingesting: {filepath.name}")
                extracted = self.extract_text(filepath)
                if extracted:
                    chunks = self.chunk_and_embed(extracted, filepath.name)
                    results[filepath.name] = chunks
                    logger.info(f"  -> {chunks} chunks indexed")
                else:
                    results[filepath.name] = 0
                    logger.warning(f"  -> No content extracted")
        return results

    def get_collection_stats(self) -> dict:
        """Get statistics about the ChromaDB collection."""
        try:
            collection = self._get_or_create_collection()
            count = collection.count()
            return {"total_chunks": count, "collection": CHROMA_COLLECTION_NAME}
        except Exception as e:
            logger.error(f"Error getting collection stats: {e}")
            return {"total_chunks": 0, "collection": CHROMA_COLLECTION_NAME}


def get_ingestor() -> DocumentIngestor:
    """Factory function to get a DocumentIngestor instance."""
    return DocumentIngestor()